"""Turn an uploaded resume file into structured, reviewable CV content.

Two stages, both deterministic and local — no external AI calls, matching
:mod:`cvbuilder.matcher` and :mod:`cvbuilder.question_extractor`:

- ``extract_text``: get plain text out of a resume file. Markdown/text is
  read directly; PDF/DOCX/PPTX go through a format-specific library.
- ``parse_resume``: split that plain text into profile paragraphs, work
  history (role/company/bullets), skills, and education, using section
  headings and blank-line paragraph breaks as the only structural signal.
  Extraction quality follows the source file's own formatting — this is
  a heuristic, not a layout-aware document parser.
"""

from __future__ import annotations

import hashlib
import io
import re
from dataclasses import asdict, dataclass, field
from pathlib import Path
from typing import Any

SUPPORTED_EXTENSIONS = {".md", ".markdown", ".txt", ".pdf", ".docx", ".pptx"}

_SUMMARY_HEADINGS = {"summary", "profile", "about", "objective", "about me"}
_EXPERIENCE_HEADINGS = {
    "experience",
    "work experience",
    "employment",
    "employment history",
    "professional experience",
    "work history",
    "career history",
}
_SKILLS_HEADINGS = {
    "skills",
    "technical skills",
    "core competencies",
    "competencies",
    "key skills",
}
_EDUCATION_HEADINGS = {
    "education",
    "education and certifications",
    "academic background",
    "qualifications",
}
_ALL_HEADINGS = (
    _SUMMARY_HEADINGS | _EXPERIENCE_HEADINGS | _SKILLS_HEADINGS | _EDUCATION_HEADINGS
)

_BULLET_RE = re.compile(r"^\s*(?:[-*•‣▪]|\d+[.)]|[a-zA-Z][.)])\s+(.*\S)\s*$")
_HEADER_SPLIT_RE = re.compile(r"\s+(?:at|@)\s+|\s*[|—–]\s*|\s+-\s+|,\s+")
_TRAILING_DATE_RE = re.compile(
    r"\s*[\(\[]?\b(19|20)\d{2}\b.*$|\s*\((?:present|current)\)\s*$", re.IGNORECASE
)
_HEADING_MAX_LEN = 40


@dataclass
class ParsedExperience:
    """One work-history entry: a role/company header plus its bullets."""

    heading: str
    role: str = ""
    company: str = ""
    bullets: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        """Return a JSON-serialisable representation."""
        return asdict(self)


@dataclass
class ParsedResume:
    """Structured content pulled out of an uploaded resume file."""

    profile: list[str] = field(default_factory=list)
    experience: list[ParsedExperience] = field(default_factory=list)
    skills: list[str] = field(default_factory=list)
    education: list[str] = field(default_factory=list)

    def to_dict(self) -> dict[str, Any]:
        """Return a JSON-serialisable representation."""
        return {
            "profile": list(self.profile),
            "experience": [item.to_dict() for item in self.experience],
            "skills": list(self.skills),
            "education": list(self.education),
        }

    def counts(self) -> dict[str, int]:
        """Return per-section item counts for a review-stage summary."""
        return {
            "profile": len(self.profile),
            "experience": len(self.experience),
            "skills": len(self.skills),
            "education": len(self.education),
        }


def extract_text(filename: str, data: bytes) -> str:
    """Extract plain text from an uploaded resume file's raw bytes.

    Args:
        filename: Original filename (used only for its extension).
        data: Raw file bytes.

    Returns:
        Extracted plain text.

    Raises:
        ValueError: Unsupported extension, or the file could not be read
            by its format's library (corrupt/password-protected/etc).
    """
    ext = Path(filename).suffix.lower()
    try:
        if ext in {".md", ".markdown", ".txt"}:
            return data.decode("utf-8", errors="replace")
        if ext == ".pdf":
            return _extract_pdf(data)
        if ext == ".docx":
            return _extract_docx(data)
        if ext == ".pptx":
            return _extract_pptx(data)
    except ValueError:
        raise
    except Exception as exc:  # noqa: BLE001 - third-party parser failures vary widely
        raise ValueError(f"could not read {ext} file: {exc}") from exc
    raise ValueError(f"unsupported file type: {ext or '(none)'}")


def _extract_pdf(data: bytes) -> str:
    """Extract text from a PDF, page by page."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = [page.extract_text() or "" for page in reader.pages]
    return "\n\n".join(pages)


def _extract_docx(data: bytes) -> str:
    """Extract text from a Word document, paragraph by paragraph."""
    from docx import Document

    document = Document(io.BytesIO(data))
    return "\n".join(paragraph.text for paragraph in document.paragraphs)


def _extract_pptx(data: bytes) -> str:
    """Extract text from a PowerPoint deck, slide by slide."""
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for slide in presentation.slides:
        lines: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                text = "".join(run.text for run in paragraph.runs).strip()
                if text:
                    lines.append(text)
        if lines:
            slides.append("\n".join(lines))
    return "\n\n".join(slides)


def parse_resume(text: str) -> ParsedResume:
    """Split resume text into profile/experience/skills/education sections.

    Section boundaries are detected from short heading-like lines that
    match a known vocabulary (see ``_ALL_HEADINGS``); everything before
    the first recognised heading is treated as the profile summary.
    """
    sections = _split_into_sections(text)
    resume = ParsedResume()
    for name, body in sections:
        if name == "profile":
            resume.profile.extend(_split_paragraphs(body))
        elif name == "experience":
            resume.experience.extend(_parse_experience(body))
        elif name == "skills":
            resume.skills.extend(_parse_skills(body))
        elif name == "education":
            resume.education.extend(_split_paragraphs(body))
    return resume


def _classify_heading(line: str) -> str | None:
    """Return the section name a heading-like line belongs to, if any."""
    stripped = line.strip().strip(":").strip()
    if not stripped or len(stripped) > _HEADING_MAX_LEN:
        return None
    key = stripped.lower()
    if key in _SUMMARY_HEADINGS:
        return "profile"
    if key in _EXPERIENCE_HEADINGS:
        return "experience"
    if key in _SKILLS_HEADINGS:
        return "skills"
    if key in _EDUCATION_HEADINGS:
        return "education"
    return None


def _split_into_sections(text: str) -> list[tuple[str, str]]:
    """Walk the document line by line, grouping text under each heading."""
    sections: list[tuple[str, list[str]]] = [("profile", [])]
    for raw_line in text.splitlines():
        heading = _classify_heading(raw_line)
        if heading is not None:
            sections.append((heading, []))
            continue
        sections[-1][1].append(raw_line)
    return [(name, "\n".join(lines)) for name, lines in sections]


def _split_paragraphs(body: str) -> list[str]:
    """Split a section body into blank-line-delimited paragraphs."""
    paragraphs: list[str] = []
    for block in re.split(r"\n\s*\n", body):
        joined = " ".join(line.strip() for line in block.splitlines() if line.strip())
        if joined:
            paragraphs.append(joined)
    return paragraphs


def _parse_skills(body: str) -> list[str]:
    """Split a skills section into individual skill entries."""
    skills: list[str] = []
    for line in body.splitlines():
        line = _strip_bullet(line.strip())
        if not line:
            continue
        if "," in line or "·" in line or "|" in line:
            parts = re.split(r"[,·|]", line)
            skills.extend(part.strip() for part in parts if part.strip())
        else:
            skills.append(line)
    seen: set[str] = set()
    unique: list[str] = []
    for skill in skills:
        key = skill.lower()
        if key in seen:
            continue
        seen.add(key)
        unique.append(skill)
    return unique


def _strip_bullet(line: str) -> str:
    """Strip a leading bullet/number marker from a line, if present."""
    match = _BULLET_RE.match(line)
    return match.group(1) if match else line


def _parse_experience(body: str) -> list[ParsedExperience]:
    """Split an experience section into role blocks on blank lines.

    Each block's first line is treated as the role/company header; the
    remaining lines become bullets (bullet markers stripped if present).
    """
    entries: list[ParsedExperience] = []
    for block in re.split(r"\n\s*\n", body):
        lines = [line.strip() for line in block.splitlines() if line.strip()]
        if not lines:
            continue
        header = lines[0]
        bullets = [_strip_bullet(line) for line in lines[1:]]
        role, company = _split_header(header)
        entries.append(
            ParsedExperience(heading=header, role=role, company=company, bullets=bullets)
        )
    return entries


def _split_header(header: str) -> tuple[str, str]:
    """Best-effort split of a role/company header line into (role, company)."""
    cleaned = _TRAILING_DATE_RE.sub("", header).strip().rstrip("-,|").strip()
    parts = [part.strip() for part in _HEADER_SPLIT_RE.split(cleaned, maxsplit=1)]
    if len(parts) == 2 and all(parts):
        return parts[0], parts[1]
    return cleaned, ""


def content_hash(content: str) -> str:
    """Return a stable SHA-256 hex digest of normalised content.

    Matches :meth:`cvbuilder.importer.SnippetImporter._hash_content` so
    imported and YAML-seeded snippets dedupe against each other.
    """
    normalised = "\n".join(line.rstrip() for line in content.strip().splitlines())
    return hashlib.sha256(normalised.encode("utf-8")).hexdigest()


def build_candidates(resume: ParsedResume) -> list[dict[str, Any]]:
    """Flatten a ``ParsedResume`` into snippet-shaped candidate dicts.

    Each dict has ``section``, ``index``, ``category``, ``heading``,
    ``company``, ``role``, ``tags``, and ``content`` — enough to preview
    (hash + dedupe-check) or create (upsert) a library snippet, without
    this module knowing anything about the database.
    """
    candidates: list[dict[str, Any]] = []
    for index, paragraph in enumerate(resume.profile):
        candidates.append(
            {
                "section": "profile",
                "index": index,
                "category": "bio",
                "heading": f"Imported summary {index + 1}",
                "company": None,
                "role": None,
                "tags": ["bio", "import"],
                "content": paragraph,
            }
        )
    for index, entry in enumerate(resume.experience):
        content = (
            "\n".join(f"- {bullet}" for bullet in entry.bullets)
            if entry.bullets
            else entry.heading
        )
        tags = ["experience", "import"]
        if entry.company:
            tags.append(entry.company.lower())
        candidates.append(
            {
                "section": "experience",
                "index": index,
                "category": "experience",
                "heading": entry.role or entry.heading,
                "company": entry.company or None,
                "role": entry.role or None,
                "tags": tags,
                "content": content,
            }
        )
    for index, skill in enumerate(resume.skills):
        candidates.append(
            {
                "section": "skills",
                "index": index,
                "category": "skill",
                "heading": skill,
                "company": None,
                "role": None,
                "tags": ["skill", "import"],
                "content": skill,
            }
        )
    for index, entry in enumerate(resume.education):
        candidates.append(
            {
                "section": "education",
                "index": index,
                "category": "part",
                "heading": "Education",
                "company": None,
                "role": None,
                "tags": ["education", "import"],
                "content": entry,
            }
        )
    return candidates
