"""Tests for resume text extraction/parsing and the resume_imports table."""

from __future__ import annotations

import io
from typing import TYPE_CHECKING

import pytest

from cvbuilder.database import SnippetDatabase
from cvbuilder.models import ResumeImport, Snippet, SnippetVariant
from cvbuilder.resume_extractor import (
    build_candidates,
    content_hash,
    extract_text,
    parse_resume,
)

if TYPE_CHECKING:
    pass

SAMPLE_RESUME = """Homer Simpson
Staff Platform Engineer

Summary
Platform and infrastructure leader with a decade of experience building
resilient cloud systems for regulated industries.

Experience
Staff Platform Engineer — Acme Corp (2021 - Present)
- Led migration of 40+ services to Kubernetes, cutting infra cost 30%
- Built the on-call incident response program from scratch

Senior Platform Engineer, Globex Inc (2017-2021)
- Owned the CI/CD platform used by 200 engineers

Skills
Kubernetes, Terraform, AWS, Python

Education
BSc Computer Science, State University, 2013
"""


class TestExtractText:
    """Format dispatch and error handling for extract_text()."""

    def test_markdown_and_text_are_read_directly(self) -> None:
        assert extract_text("resume.md", b"# Hello") == "# Hello"
        assert extract_text("resume.txt", b"Hello") == "Hello"

    def test_docx_extraction(self) -> None:
        from docx import Document

        document = Document()
        document.add_paragraph("Homer Simpson")
        document.add_paragraph("Staff Engineer at Acme")
        buf = io.BytesIO()
        document.save(buf)
        text = extract_text("resume.docx", buf.getvalue())
        assert "Homer Simpson" in text
        assert "Staff Engineer at Acme" in text

    def test_pptx_extraction(self) -> None:
        from pptx import Presentation

        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = "Skills"
        buf = io.BytesIO()
        presentation.save(buf)
        text = extract_text("resume.pptx", buf.getvalue())
        assert "Skills" in text

    def test_unsupported_extension_raises(self) -> None:
        with pytest.raises(ValueError, match="unsupported file type"):
            extract_text("resume.exe", b"whatever")

    def test_corrupt_pdf_raises_value_error(self) -> None:
        with pytest.raises(ValueError, match="could not read .pdf file"):
            extract_text("resume.pdf", b"not a real pdf")


class TestParseResume:
    """Heuristic section splitting on the deterministic extractor."""

    def test_splits_profile_experience_skills_education(self) -> None:
        resume = parse_resume(SAMPLE_RESUME)
        assert len(resume.profile) == 2
        assert "decade of experience" in resume.profile[1]
        assert len(resume.experience) == 2
        assert resume.skills == ["Kubernetes", "Terraform", "AWS", "Python"]
        assert resume.education == ["BSc Computer Science, State University, 2013"]

    def test_experience_role_and_company_split(self) -> None:
        resume = parse_resume(SAMPLE_RESUME)
        first = resume.experience[0]
        assert first.role == "Staff Platform Engineer"
        assert first.company == "Acme Corp"
        assert len(first.bullets) == 2
        assert first.bullets[0].startswith("Led migration")

    def test_no_recognised_headings_puts_everything_in_profile(self) -> None:
        resume = parse_resume("Just a couple of\nplain lines with no headings.")
        assert resume.profile == ["Just a couple of plain lines with no headings."]
        assert resume.experience == []
        assert resume.skills == []
        assert resume.education == []

    def test_empty_text_produces_empty_resume(self) -> None:
        resume = parse_resume("")
        assert resume.counts() == {
            "profile": 0,
            "experience": 0,
            "skills": 0,
            "education": 0,
        }

    def test_selected_experience_after_skills_is_not_lost(self) -> None:
        """PDF layouts often emit Technical Skills before Selected Experience."""
        text = """
Ramon Brooker
DevOps Evangelist

Technical Skills
Kubernetes, Terraform, Python

Business Founder / Entrepreneur
AEO3 Ltd.
Founder | December 2023 - Present | Ottawa, Ontario
- Defined company goals and strategic direction

Selected Experience
Accenture
Application Development Manager | May 2020 Oct 2022 | Ottawa, Ontario
- Led mixed teams from across the globe

Imagine Communications Canada, Ltd
Senior Software Engineer, R&D | Aug 2014 May 2020 | Waterloo, Ontario
- Designed and deployed Playout Platforms

Training and Education
BSc Something, Some University
""".strip()
        resume = parse_resume(text)
        assert resume.experience, "expected experience entries to be detected"
        assert len(resume.experience) >= 2
        joined = " ".join(
            f"{entry.heading} {entry.role} {entry.company}"
            for entry in resume.experience
        )
        assert "Accenture" in joined or "Application Development Manager" in joined
        assert "Imagine" in joined or "Senior Software Engineer" in joined
        assert "Kubernetes" in resume.skills or "Python" in resume.skills
        assert resume.education, "Training and Education should map to education"


class TestBuildCandidates:
    """Flattening a ParsedResume into snippet-shaped candidate dicts."""

    def test_candidate_counts_match_parsed_resume(self) -> None:
        resume = parse_resume(SAMPLE_RESUME)
        candidates = build_candidates(resume)
        counts = resume.counts()
        assert len(candidates) == sum(counts.values())
        sections = {c["section"] for c in candidates}
        assert sections == {"profile", "experience", "skills", "education"}

    def test_experience_candidate_content_is_bulleted(self) -> None:
        resume = parse_resume(SAMPLE_RESUME)
        candidates = build_candidates(resume)
        experience = next(c for c in candidates if c["section"] == "experience")
        assert experience["company"] == "Acme Corp"
        assert experience["content"].startswith("- Led migration")


class TestResumeImportsTable:
    """CRUD for the resume_imports table."""

    def test_create_list_get_delete(self, snippet_db: SnippetDatabase) -> None:
        import_id = snippet_db.create_resume_import(
            ResumeImport(
                filename="resume.pdf",
                file_type="pdf",
                stored_path="abc123__resume.pdf",
                snippet_count=7,
            )
        )
        imports = snippet_db.list_resume_imports()
        assert len(imports) == 1
        assert imports[0].id == import_id
        assert imports[0].filename == "resume.pdf"
        assert imports[0].snippet_count == 7

        fetched = snippet_db.get_resume_import(import_id)
        assert fetched is not None
        assert fetched.stored_path == "abc123__resume.pdf"

        assert snippet_db.delete_resume_import(import_id) is True
        assert snippet_db.list_resume_imports() == []

    def test_existing_content_hashes(self, snippet_db: SnippetDatabase) -> None:
        snippet = Snippet(
            category="skill",
            heading="Kubernetes",
            source_path="cv/web/data.yaml#skills.technical[0]",
            content_hash=content_hash("Kubernetes"),
        )
        variant = SnippetVariant(detail_level="standard", content="Kubernetes")
        snippet_db.upsert_by_source(snippet, variant)

        hashes = [content_hash("Kubernetes"), content_hash("Terraform")]
        existing = snippet_db.existing_content_hashes(hashes)
        assert existing == {content_hash("Kubernetes")}
